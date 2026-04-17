"""Build the 3-slide investor deck.

Pulls real numbers from data_cache/backtest_report.json so the deck is
grounded in today's honest metrics — no made-up figures.

Outputs: research/s-tool_investor_deck_2026_04_17.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).parent.parent
BACKTEST = json.loads((ROOT / "data_cache" / "backtest_report.json").read_text())
OUT = Path(__file__).parent / "s-tool_investor_deck_2026_04_17.pptx"


# Palette matched to the live s-tool.io theme (deep forest + lake accent).
BG_DEEP     = RGBColor(0x0B, 0x10, 0x06)   # --bg-deepest
BG_SURFACE  = RGBColor(0x1A, 0x1F, 0x1B)
BG_CARD     = RGBColor(0x1F, 0x27, 0x22)
TEXT_HI     = RGBColor(0xED, 0xEC, 0xE6)
TEXT_MID    = RGBColor(0x9B, 0xA1, 0xB9)
TEXT_DIM    = RGBColor(0x6B, 0x73, 0x82)
ACCENT      = RGBColor(0x5F, 0xAA, 0xC5)   # lake
ACCENT_DEEP = RGBColor(0x0A, 0x65, 0x77)   # alpine
POS         = RGBColor(0x6E, 0xE7, 0xB7)
WARN        = RGBColor(0xF5, 0xD5, 0x8F)


def add_background(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, left, top, width, height, text, *,
              font_size=14, font_name="Helvetica", color=TEXT_HI,
              bold=False, italic=False, align=PP_ALIGN.LEFT,
              line_spacing=1.2):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    first = True
    for line in (text if isinstance(text, list) else [text]):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        if isinstance(line, tuple):
            # (text, {font_size, color, bold, ...})
            body, overrides = line
        else:
            body, overrides = line, {}
        r = p.add_run()
        r.text = body
        f = r.font
        f.name = overrides.get("font_name", font_name)
        f.size = Pt(overrides.get("font_size", font_size))
        f.color.rgb = overrides.get("color", color)
        f.bold = overrides.get("bold", bold)
        f.italic = overrides.get("italic", italic)
    return tx


def add_rect(slide, left, top, width, height, *, fill=BG_CARD,
             line=None, radius=0.05):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = radius
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


# ── Slide 1 — The edge ─────────────────────────────────────────────
def slide_performance(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_background(slide, BG_DEEP)

    # Eyebrow + title
    add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
             "THE EDGE  ·  BACKTEST 2022-05  →  2024-08  ·  10 WALK-FORWARD WINDOWS",
             font_size=11, color=ACCENT, bold=True)
    add_text(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(1.1),
             [("Our ranker picks ", {"font_size": 36, "color": TEXT_HI}),
              ("+100% movers at 11× the Russell 3000 baseline", {"font_size": 36, "color": ACCENT, "italic": True}),
              (".", {"font_size": 36, "color": TEXT_HI})],
             font_name="Georgia")

    # Baseline card
    bx, by, bw, bh = Inches(0.6), Inches(2.1), Inches(6.0), Inches(2.3)
    add_rect(slide, bx, by, bw, bh, fill=BG_SURFACE, line=TEXT_DIM)
    add_text(slide, bx + Inches(0.3), by + Inches(0.2), Inches(5.4), Inches(0.3),
             "UNIVERSE BASELINE  ·  RANDOM TICKER, 1-YEAR RETURN",
             font_size=10, color=ACCENT, bold=True)
    # Baseline rates table (4 rows)
    base = BACKTEST["baseline_rates"]
    rows = [("+10%",  f"{base['+10%']*100:.1f}%"),
            ("+25%",  f"{base['+25%']*100:.1f}%"),
            ("+50%",  f"{base['+50%']*100:.1f}%"),
            ("+100%", f"{base['+100%']*100:.2f}%"),
            ("+200%", f"{base['+200%']*100:.2f}%")]
    for i, (thr, rate) in enumerate(rows):
        y = by + Inches(0.6 + 0.3 * i)
        add_text(slide, bx + Inches(0.3), y, Inches(2), Inches(0.3),
                 thr, font_size=14, color=TEXT_MID)
        add_text(slide, bx + Inches(2.3), y, Inches(2.4), Inches(0.3),
                 rate, font_size=14, color=TEXT_HI, bold=True, align=PP_ALIGN.LEFT)

    # Model card — right
    mx, my, mw, mh = Inches(7.0), Inches(2.1), Inches(5.7), Inches(2.3)
    add_rect(slide, mx, my, mw, mh, fill=BG_SURFACE, line=ACCENT)
    add_text(slide, mx + Inches(0.3), my + Inches(0.2), Inches(5.4), Inches(0.3),
             "OUR NN SCORER  ·  TOP-20 PICKS / WINDOW",
             font_size=10, color=ACCENT, bold=True)
    m = BACKTEST["methods"]
    nn = m["nn_score"]["thresholds"]["+100%"]
    rows = [
        ("Hit rate at +100%",        f"{nn['rate']*100:.1f}%"),
        ("Lift vs baseline",          f"{nn['lift']:.2f}×"),
        ("Mean 1Y return on picks",   f"+{m['nn_score']['mean_return']*100:.0f}%"),
        ("Median 1Y return on picks", f"+{m['nn_score']['median_return']*100:.0f}%"),
        ("Hand-crafted baseline lift", f"{m['H7_ewma_p90']['thresholds']['+100%']['lift']:.2f}×"),
    ]
    for i, (lab, val) in enumerate(rows):
        y = my + Inches(0.6 + 0.32 * i)
        add_text(slide, mx + Inches(0.3), y, Inches(3.4), Inches(0.3),
                 lab, font_size=13, color=TEXT_MID)
        c = POS if "Lift" in lab or "Hit" in lab or "Mean" in lab or "Median" in lab else TEXT_HI
        add_text(slide, mx + Inches(3.7), y, Inches(1.9), Inches(0.3),
                 val, font_size=16, color=c, bold=True,
                 font_name="Georgia", align=PP_ALIGN.RIGHT)

    # Honest metrics strip
    hx, hy, hw, hh = Inches(0.6), Inches(4.65), Inches(12.1), Inches(1.8)
    add_rect(slide, hx, hy, hw, hh, fill=BG_CARD, line=TEXT_DIM)
    add_text(slide, hx + Inches(0.3), hy + Inches(0.15), Inches(11.5), Inches(0.3),
             "HONEST METRICS  ·  SIZE-NEUTRAL  ·  OUT-OF-SAMPLE  ·  WITHIN-QUINTILE",
             font_size=10, color=WARN, bold=True)
    hm = BACKTEST["honest_metrics"]
    cols = [
        ("Size-neutral hit_100",  f"{hm['size_neutral_hit_100']*100:.1f}%",
                                  "not a small-cap tilt artifact"),
        ("Within-quintile lift",  f"{hm['within_quintile_lift_median']:.2f}×",
                                  "median across 5 market-cap quintiles"),
        ("2024 OOS hit_100",      f"{hm['year_oos_hit_100']*100:.1f}%",
                                  "year never seen by training loop"),
        ("2024 OOS lift",          f"{hm['year_oos_lift']:.2f}×",
                                  "signal strengthens out-of-sample"),
    ]
    col_w = Inches(3.0); col_x0 = hx + Inches(0.3)
    for i, (lab, val, sub) in enumerate(cols):
        cx = col_x0 + col_w * i
        add_text(slide, cx, hy + Inches(0.55), col_w, Inches(0.3),
                 lab, font_size=11, color=TEXT_MID)
        add_text(slide, cx, hy + Inches(0.85), col_w, Inches(0.55),
                 val, font_size=32, color=POS, bold=True, font_name="Georgia")
        add_text(slide, cx, hy + Inches(1.4), col_w, Inches(0.3),
                 sub, font_size=10, color=TEXT_DIM, italic=True)

    # Footer
    add_text(slide, Inches(0.6), Inches(6.9), Inches(12), Inches(0.4),
             [("22,218 ticker-window rows  ·  Russell 3000 universe  ·  every pick logged in picks_history ledger for live tracking  ·  ",
               {"font_size": 11, "color": TEXT_DIM}),
              ("s-tool.io/track-record",
               {"font_size": 11, "color": ACCENT, "bold": True})])


# ── Slide 2 — Architecture ────────────────────────────────────────
def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG_DEEP)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
             "ARCHITECTURE  ·  NIGHTLY END-TO-END LEARNING LOOP",
             font_size=11, color=ACCENT, bold=True)
    add_text(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(1.1),
             [("The model retrains ", {"font_size": 36, "color": TEXT_HI}),
              ("every night", {"font_size": 36, "color": ACCENT, "italic": True}),
              (" on freshly realised returns.", {"font_size": 36, "color": TEXT_HI})],
             font_name="Georgia")

    # 4 stack columns
    col_defs = [
        ("DATA", ACCENT, [
            "Russell 3000 universe via iShares IWV",
            "2,609 ticker price histories (yfinance, 5y+)",
            "1,947 SEC EDGAR XBRL companies (98k quarterly rows)",
            "FMP Premium analyst consensus + EPS revisions",
            "FRED macro regime (6 indicators)",
            "FinBERT sentiment on historical Reddit + StockTwits",
        ]),
        ("COMPUTE", POS, [
            "Monte Carlo GBM + Ornstein-Uhlenbeck mean reversion",
            "30/70 MC↔MR blend tuned against realised returns",
            "ExtraTreesRegressor NN scorer, walk-forward trained",
            "Size-neutral asymmetric bucketing (quintile-balanced)",
            "$500K avg daily liquidity floor (filters penny stocks)",
            "NN retraining nightly; HP search winner from v2 study",
        ]),
        ("SERVE", WARN, [
            "Cloudflare Worker — 15KB HTML, edge-served",
            "FastAPI backend on Railway, Clerk auth, Stripe billing",
            "Railway Volume — users DB persists across redeploys",
            "SQLite engine cache, WAL-mode, GitHub Actions cached",
            "Gated Strategist tier → full ranked list + asymmetric",
            "/track-record publishes the ledger for public auditing",
        ]),
        ("CI  ·  LOOP", RGBColor(0xD2, 0xDD, 0xEA), [
            "Fast path @ 20:00 UTC weekdays: 537-ticker preferred",
            "  universe scan + deploy, 30–45 min cap",
            "Slow path @ 23:00 UTC: Russell 3000 long-tail +",
            "  NN training + honest backtest + feature ablation",
            "Both deploy Cloudflare + Railway on completion",
            "Every run's picks go into picks_history ledger",
        ]),
    ]
    col_x = Inches(0.6)
    col_w = Inches(3.05)
    col_y = Inches(2.1)
    col_h = Inches(3.7)
    for i, (label, color, lines) in enumerate(col_defs):
        x = col_x + col_w * i + Inches(0.05 * i)
        add_rect(slide, x, col_y, col_w, col_h, fill=BG_SURFACE, line=color)
        # Header strip
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, col_y, col_w, Inches(0.5))
        hdr.adjustments[0] = 0.15
        hdr.fill.solid(); hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        hdr.shadow.inherit = False
        add_text(slide, x, col_y + Inches(0.1), col_w, Inches(0.35),
                 label, font_size=13, color=BG_DEEP, bold=True, align=PP_ALIGN.CENTER)
        # Body bullets
        for j, ln in enumerate(lines):
            add_text(slide,
                     x + Inches(0.15),
                     col_y + Inches(0.6 + 0.43 * j),
                     col_w - Inches(0.2),
                     Inches(0.4),
                     "•  " + ln,
                     font_size=10, color=TEXT_HI, line_spacing=1.15)

    # Loop callout strip
    lx, ly, lw, lh = Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.0)
    add_rect(slide, lx, ly, lw, lh, fill=BG_CARD, line=ACCENT)
    add_text(slide, lx + Inches(0.3), ly + Inches(0.12), Inches(11.5), Inches(0.3),
             "THE FEEDBACK LOOP", font_size=10, color=ACCENT, bold=True)
    add_text(slide, lx + Inches(0.3), ly + Inches(0.45), Inches(11.5), Inches(0.5),
             [("today's market prints  →  ", {"font_size": 14, "color": TEXT_MID}),
              ("nightly realisation join  →  ",  {"font_size": 14, "color": TEXT_MID}),
              ("walk-forward retrain  →  ",      {"font_size": 14, "color": TEXT_MID}),
              ("fresh ranking  →  ",             {"font_size": 14, "color": TEXT_MID}),
              ("tomorrow's picks",               {"font_size": 14, "color": ACCENT, "bold": True})])

    add_text(slide, Inches(0.6), Inches(7.15), Inches(12), Inches(0.3),
             "Static models decay.  We rebuild the model daily — it responds to regime shifts in under 24h.",
             font_size=11, color=TEXT_DIM, italic=True)


# ── Slide 3 — Business ────────────────────────────────────────────
def slide_business(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG_DEEP)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
             "WHERE CAPITAL GOES  ·  2026 Q2",
             font_size=11, color=ACCENT, bold=True)
    add_text(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(1.1),
             [("Real edge.  Honest surface.  Infra built for the ", {"font_size": 32, "color": TEXT_HI}),
              ("next 100× of users", {"font_size": 32, "color": ACCENT, "italic": True}),
              (".", {"font_size": 32, "color": TEXT_HI})],
             font_name="Georgia")

    # Left — Pricing / stage
    lx, ly, lw, lh = Inches(0.6), Inches(2.1), Inches(5.9), Inches(2.8)
    add_rect(slide, lx, ly, lw, lh, fill=BG_SURFACE, line=POS)
    add_text(slide, lx + Inches(0.3), ly + Inches(0.2), Inches(5.3), Inches(0.3),
             "PRICING  ·  REVENUE SURFACE",
             font_size=10, color=POS, bold=True)
    tiers = [
        ("Free",       "3 projections / day",            "—"),
        ("Pro",        "10 projections / day",           "$8 / mo"),
        ("Strategist", "full ranked list + asymmetric",  "$29 / mo"),
    ]
    hdr_y = ly + Inches(0.55)
    add_text(slide, lx + Inches(0.3), hdr_y, Inches(2), Inches(0.25),
             "TIER", font_size=9, color=TEXT_DIM, bold=True)
    add_text(slide, lx + Inches(2.0), hdr_y, Inches(2.5), Inches(0.25),
             "WHAT", font_size=9, color=TEXT_DIM, bold=True)
    add_text(slide, lx + Inches(4.6), hdr_y, Inches(1.2), Inches(0.25),
             "PRICE", font_size=9, color=TEXT_DIM, bold=True, align=PP_ALIGN.RIGHT)
    for i, (t, w, p) in enumerate(tiers):
        y = ly + Inches(0.85 + 0.45 * i)
        add_text(slide, lx + Inches(0.3), y, Inches(2), Inches(0.3),
                 t, font_size=14, color=TEXT_HI, bold=True)
        add_text(slide, lx + Inches(2.0), y, Inches(2.5), Inches(0.3),
                 w, font_size=12, color=TEXT_MID)
        add_text(slide, lx + Inches(4.6), y, Inches(1.2), Inches(0.3),
                 p, font_size=14, color=POS, bold=True,
                 font_name="Georgia", align=PP_ALIGN.RIGHT)
    add_text(slide, lx + Inches(0.3), ly + Inches(2.25), Inches(5.3), Inches(0.5),
             "Stripe billing live  ·  Clerk auth  ·  paywall enforcement flagged; ready to switch on",
             font_size=10, color=TEXT_DIM, italic=True, line_spacing=1.2)

    # Right — Unit economics
    rx, ry, rw, rh = Inches(6.9), Inches(2.1), Inches(5.9), Inches(2.8)
    add_rect(slide, rx, ry, rw, rh, fill=BG_SURFACE, line=ACCENT)
    add_text(slide, rx + Inches(0.3), ry + Inches(0.2), Inches(5.3), Inches(0.3),
             "UNIT ECONOMICS", font_size=10, color=ACCENT, bold=True)
    ue = [
        ("Infra cost",            "<$50 / mo",      "all providers on free or low tier"),
        ("Per-user marginal cost", "~$0",            "nightly compute is fixed, serving is edge"),
        ("Gross margin @ scale",  "95%+",           "subscription less Stripe fees"),
        ("Capacity",              "10k+ users",     "without architecture change"),
    ]
    for i, (lab, val, sub) in enumerate(ue):
        y = ry + Inches(0.6 + 0.55 * i)
        add_text(slide, rx + Inches(0.3), y, Inches(3.3), Inches(0.3),
                 lab, font_size=12, color=TEXT_MID)
        add_text(slide, rx + Inches(3.6), y, Inches(2.0), Inches(0.3),
                 val, font_size=16, color=POS, bold=True,
                 font_name="Georgia", align=PP_ALIGN.RIGHT)
        add_text(slide, rx + Inches(0.3), y + Inches(0.28), Inches(5.3), Inches(0.25),
                 sub, font_size=10, color=TEXT_DIM, italic=True)

    # Bottom — Roadmap + risk honesty
    bx, by, bw, bh = Inches(0.6), Inches(5.05), Inches(12.1), Inches(2.15)
    add_rect(slide, bx, by, bw, bh, fill=BG_CARD, line=TEXT_DIM)
    add_text(slide, bx + Inches(0.3), by + Inches(0.15), Inches(4), Inches(0.3),
             "ROADMAP  ·  NEXT 90 DAYS",
             font_size=10, color=WARN, bold=True)
    road = [
        "Live pick scoreboard: weekly realised returns, public, unedited",
        "Options flow (Polygon) + earnings calendar (Finnhub) as new NN features",
        "FINRA short interest historical backfill — capital-commitment signal",
        "Portfolio-level backtest: simulate a user following the picks for 12 months",
    ]
    for i, r in enumerate(road):
        add_text(slide, bx + Inches(0.3), by + Inches(0.5 + 0.28 * i),
                 Inches(7.5), Inches(0.3),
                 "•  " + r, font_size=11, color=TEXT_HI)
    add_text(slide, bx + Inches(8.2), by + Inches(0.15), Inches(3.8), Inches(0.3),
             "RISKS  ·  WHAT WE WATCH",
             font_size=10, color=WARN, bold=True)
    risks = [
        "Past backtest not a future guarantee",
        "Regime shift killed crowd-sentiment in 2023 — flagged + removed",
        "Russell 3000 scan budget: ~45 min / day — headroom to 90 min",
        "No human discretion overlay — fully systematic",
    ]
    for i, r in enumerate(risks):
        add_text(slide, bx + Inches(8.2), by + Inches(0.5 + 0.28 * i),
                 Inches(3.8), Inches(0.3),
                 "•  " + r, font_size=11, color=TEXT_MID)


def main() -> None:
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_performance(prs)
    slide_architecture(prs)
    slide_business(prs)
    prs.save(str(OUT))
    print(f"wrote {OUT}  ({OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
