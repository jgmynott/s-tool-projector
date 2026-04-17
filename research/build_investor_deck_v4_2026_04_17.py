"""Investor deck v4 — consulting-grade spacing and alignment.

v3 fixed vocabulary but had text floating inside cards, inconsistent
padding, mixed bullet markers, and pricing rows with off-kilter
alignment. v4 enforces a strict layout system:

  1. Every card has CARD_PAD_IN = 0.35 on all sides.
  2. Every card header is a fixed 0.5in tall coloured strip at top.
  3. Every "hero stat" element uses MSO_ANCHOR.MIDDLE so the text
     visually centres inside its bounding box rather than top-stacking
     and leaving trailing whitespace.
  4. Native PowerPoint tables handle all tabular data (pricing, unit
     economics) — rows auto-align and cell padding is consistent.
  5. Bullets are the same "•" glyph everywhere with a fixed indent.
  6. Per-slide thumbnails regenerated for visual QC before commit.
"""

from __future__ import annotations

import json
import shutil
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).parent.parent
BACKTEST = json.loads((ROOT / "data_cache" / "backtest_report.json").read_text())
OUT = Path(__file__).parent / "s-tool_investor_deck_v4_2026_04_17.pptx"

# ── Strict layout constants ──
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.5
GUTTER = 0.3
CARD_PAD = 0.35           # inner padding on every card — FIXED
HDR_H = 0.5               # height of coloured card-header strip
TITLE_Y = 0.4
TITLE_H = 1.1
BODY_Y = 1.9              # where body content starts below title

# ── Palette ──
BG_DEEP    = RGBColor(0x0B, 0x10, 0x06)
BG_SURFACE = RGBColor(0x1A, 0x1F, 0x1B)
BG_CARD    = RGBColor(0x1F, 0x27, 0x22)
TEXT_HI    = RGBColor(0xED, 0xEC, 0xE6)
TEXT_MID   = RGBColor(0x9B, 0xA1, 0xB9)
TEXT_DIM   = RGBColor(0x6B, 0x73, 0x82)
ACCENT     = RGBColor(0x5F, 0xAA, 0xC5)
POS        = RGBColor(0x6E, 0xE7, 0xB7)
WARN       = RGBColor(0xF5, 0xD5, 0x8F)


# ── Primitive helpers ──
def add_bg(slide):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  Inches(SLIDE_W), Inches(SLIDE_H))
    shp.fill.solid(); shp.fill.fore_color.rgb = BG_DEEP
    shp.line.fill.background(); shp.shadow.inherit = False


def add_card(slide, l, t, w, h, *, fill=BG_SURFACE, border=None,
              header_label=None, header_color=ACCENT):
    """Rounded card with optional fixed-height coloured header strip."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(l), Inches(t), Inches(w), Inches(h))
    shp.adjustments[0] = 0.04
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border; shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if header_label:
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(l), Inches(t),
                                     Inches(w), Inches(HDR_H))
        hdr.adjustments[0] = 0.18
        hdr.fill.solid(); hdr.fill.fore_color.rgb = header_color
        hdr.line.fill.background(); hdr.shadow.inherit = False
        add_text(slide, l, t + 0.12, w, HDR_H - 0.15,
                 header_label, size=13, color=BG_DEEP, bold=True,
                 align=PP_ALIGN.CENTER)


def add_text(slide, l, t, w, h, text, *, size=14, color=TEXT_HI,
              bold=False, italic=False, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, font="Helvetica Neue",
              line_spacing=1.2):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font; f.size = Pt(size); f.color.rgb = color
    f.bold = bold; f.italic = italic


def add_title(slide, eyebrow, title_text, *, accent_phrase=None,
               size=30):
    add_text(slide, MARGIN, TITLE_Y,
             SLIDE_W - 2 * MARGIN, 0.3,
             eyebrow, size=10, color=ACCENT, bold=True)
    tx = slide.shapes.add_textbox(Inches(MARGIN), Inches(TITLE_Y + 0.4),
                                   Inches(SLIDE_W - 2 * MARGIN),
                                   Inches(TITLE_H))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.line_spacing = 1.05

    if accent_phrase and accent_phrase in title_text:
        before, _, after = title_text.partition(accent_phrase)
        segs = [(before, False), (accent_phrase, True), (after, False)]
    else:
        segs = [(title_text, False)]
    for txt, is_accent in segs:
        if not txt:
            continue
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = "Georgia"; f.size = Pt(size)
        f.color.rgb = ACCENT if is_accent else TEXT_HI
        f.italic = is_accent


def add_page_number(slide, n):
    add_text(slide, SLIDE_W - 0.7, SLIDE_H - 0.32, 0.5, 0.2,
             str(n), size=9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


def add_table(slide, l, t, w, h, rows, *, header_row=True,
               col_widths=None, row_heights=None,
               font="Helvetica Neue", body_size=13, header_size=10):
    """Native PowerPoint table. rows is list of list of (text, opts)."""
    n_rows = len(rows); n_cols = len(rows[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                       Inches(l), Inches(t),
                                       Inches(w), Inches(h))
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    if row_heights:
        for i, rh in enumerate(row_heights):
            tbl.rows[i].height = Inches(rh)
    for ri, row in enumerate(rows):
        for ci, cell_spec in enumerate(row):
            if isinstance(cell_spec, tuple):
                text, opts = cell_spec
            else:
                text, opts = cell_spec, {}
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_SURFACE if ri > 0 or not header_row else BG_CARD
            cell.margin_left = Inches(0.15)
            cell.margin_right = Inches(0.15)
            cell.margin_top = Inches(0.06)
            cell.margin_bottom = Inches(0.06)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = opts.get("align", PP_ALIGN.LEFT)
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = font
            is_header = (ri == 0 and header_row)
            f.size = Pt(opts.get("size", header_size if is_header else body_size))
            f.color.rgb = opts.get("color",
                                    TEXT_DIM if is_header else TEXT_HI)
            f.bold = opts.get("bold", is_header)
    return tbl


# ── Slide 1 — The Edge ────────────────────────────────────────────
def slide_one(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s,
              "THE EDGE  ·  BACKTESTED MAY 2022 THROUGH AUGUST 2024",
              "Our picks double at 11× the rate of random stock picking.",
              accent_phrase="11× the rate of random stock picking")

    # Two comparison cards — identical structure, centred content.
    top_y = 2.4
    top_h = 2.5
    col_w = (SLIDE_W - 2 * MARGIN - GUTTER) / 2
    lx = MARGIN
    rx = MARGIN + col_w + GUTTER

    def comparison_card(x, tone, header, big, bigcolor, sub, tail):
        add_card(s, x, top_y, col_w, top_h,
                 fill=BG_SURFACE, border=tone,
                 header_label=header, header_color=tone)
        # Big number — vertical-centred in its own box
        add_text(s, x + CARD_PAD, top_y + HDR_H + 0.15,
                 col_w - 2 * CARD_PAD, 1.05,
                 big, size=54, color=bigcolor, bold=True,
                 font="Georgia", align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Supporting phrase
        add_text(s, x + CARD_PAD, top_y + HDR_H + 1.25,
                 col_w - 2 * CARD_PAD, 0.3,
                 sub, size=14, color=TEXT_MID, italic=True,
                 align=PP_ALIGN.CENTER)
        # Citation
        add_text(s, x + CARD_PAD, top_y + HDR_H + 1.58,
                 col_w - 2 * CARD_PAD, 0.25,
                 tail, size=10, color=TEXT_DIM,
                 align=PP_ALIGN.CENTER)

    comparison_card(
        lx, TEXT_DIM,
        "RANDOM STOCK FROM THE RUSSELL 3000",
        "1 in 20", TEXT_HI,
        "chance of doubling over the next year",
        "4.95% base rate across 22,218 observations")

    comparison_card(
        rx, POS,
        "TOP 20 PICKS FROM OUR MODEL",
        "11 in 20", POS,
        "doubled in the following year",
        "55.0% realized rate on top-ranked picks")

    # Honest-check strip — 3 uniform cards
    chk_y = 5.15
    chk_h = 1.7
    chks = [
        (POS,
         "WORKS ACROSS EVERY COMPANY SIZE",
         "8.45× lift",
         "Across the smallest to largest\nRussell 3000 market-cap buckets"),
        (POS,
         "TESTED ON UNSEEN DATA",
         "68% doubled",
         "Out-of-Sample (OOS) 2024 test:\nthe signal strengthened, not weakened"),
        (POS,
         "BEATS SIMPLE RULES BY 5×",
         "11.1× vs 2.0×",
         "Our model's lift vs the best\nhand-crafted rule we tried"),
    ]
    c_total = SLIDE_W - 2 * MARGIN
    chk_w = (c_total - 2 * GUTTER) / 3
    for i, (tone, hdr, big, sub) in enumerate(chks):
        x = MARGIN + (chk_w + GUTTER) * i
        add_card(s, x, chk_y, chk_w, chk_h,
                 fill=BG_CARD, border=tone,
                 header_label=hdr, header_color=tone)
        add_text(s, x + CARD_PAD, chk_y + HDR_H + 0.15,
                 chk_w - 2 * CARD_PAD, 0.5,
                 big, size=26, color=POS, bold=True,
                 font="Georgia", align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + CARD_PAD, chk_y + HDR_H + 0.75,
                 chk_w - 2 * CARD_PAD, 0.45,
                 sub, size=11, color=TEXT_MID, italic=True,
                 align=PP_ALIGN.CENTER, line_spacing=1.3,
                 anchor=MSO_ANCHOR.TOP)

    add_text(s, MARGIN, SLIDE_H - 0.42, SLIDE_W - 2 * MARGIN, 0.25,
             "Every pick logged publicly for live verification at s-tool.io/track-record",
             size=10, color=TEXT_DIM, italic=True)
    add_page_number(s, 1)


# ── Slide 2 — How it works ────────────────────────────────────────
def slide_two(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "HOW IT WORKS",
              "The model rebuilds itself every night.",
              accent_phrase="every night")

    col_y = 2.4
    col_h = 3.5
    cols = 4
    total_w = SLIDE_W - 2 * MARGIN
    col_w = (total_w - GUTTER * (cols - 1)) / cols
    concepts = [
        (ACCENT, "WIDE UNIVERSE",   "~3,000 companies",
         "The Russell 3000 — the broadest US stock "
         "index. We rank the whole market, never cherry-pick."),
        (POS,    "RICH DATA",       "Public data only",
         "Prices, SEC filings, analyst estimates, and "
         "macroeconomic indicators from the Federal Reserve."),
        (WARN,   "MACHINE LEARNING", "Trained on history",
         "A machine-learning model learns which combinations of "
         "signals preceded real +100% moves in the past."),
        (RGBColor(0xD2, 0xDD, 0xEA), "DAILY RETRAIN", "Under 24 hours",
         "Every night the model sees yesterday's real market "
         "prints. No static algorithm decaying in a drawer."),
    ]
    for i, (tone, label, headline, body) in enumerate(concepts):
        x = MARGIN + (col_w + GUTTER) * i
        add_card(s, x, col_y, col_w, col_h,
                 fill=BG_SURFACE, border=tone,
                 header_label=label, header_color=tone)
        # Headline — vertically centred in the upper half
        add_text(s, x + CARD_PAD, col_y + HDR_H + 0.1,
                 col_w - 2 * CARD_PAD, 0.85,
                 headline, size=20, color=TEXT_HI, bold=True,
                 font="Georgia", align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Body — vertically centred in the lower half
        add_text(s, x + CARD_PAD, col_y + HDR_H + 1.0,
                 col_w - 2 * CARD_PAD, col_h - HDR_H - 1.2,
                 body, size=12, color=TEXT_MID,
                 align=PP_ALIGN.CENTER, line_spacing=1.4,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Loop strip
    loop_y = 6.15
    loop_h = 0.85
    add_card(s, MARGIN, loop_y, SLIDE_W - 2 * MARGIN, loop_h,
             fill=BG_CARD, border=ACCENT)
    add_text(s, MARGIN + CARD_PAD, loop_y + 0.1,
             SLIDE_W - 2 * MARGIN - 2 * CARD_PAD, 0.25,
             "THE FEEDBACK LOOP",
             size=10, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN + CARD_PAD, loop_y + 0.38,
             SLIDE_W - 2 * MARGIN - 2 * CARD_PAD, 0.4,
             "yesterday's market prints  →  realized outcomes  →  "
             "model updates  →  today's ranked picks",
             size=14, color=TEXT_HI, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_page_number(s, 2)


# ── Slide 3 — Business ────────────────────────────────────────────
def slide_three(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "WHERE THE CAPITAL GOES  ·  2026 Q2",
              "Real edge. Honest surface. Infrastructure built for scale.",
              accent_phrase="Infrastructure built for scale",
              size=28)

    col_w = (SLIDE_W - 2 * MARGIN - GUTTER) / 2
    lx = MARGIN
    rx = MARGIN + col_w + GUTTER

    # Pricing — native table inside a card
    top_y = 2.5
    top_h = 2.3
    add_card(s, lx, top_y, col_w, top_h,
             fill=BG_SURFACE, border=POS,
             header_label="PRICING  ·  LIVE TODAY", header_color=POS)
    pricing_rows = [
        [("Tier", {}), ("What you get", {}), ("Price", {"align": PP_ALIGN.RIGHT})],
        [("Free",       {"bold": True, "color": TEXT_HI, "size": 14}),
         ("3 projections per day",     {"size": 12, "color": TEXT_MID}),
         ("Free",                      {"bold": True, "color": POS, "size": 14, "align": PP_ALIGN.RIGHT})],
        [("Pro",        {"bold": True, "color": TEXT_HI, "size": 14}),
         ("10 projections per day",    {"size": 12, "color": TEXT_MID}),
         ("$8 / month",                {"bold": True, "color": POS, "size": 14, "align": PP_ALIGN.RIGHT})],
        [("Strategist", {"bold": True, "color": TEXT_HI, "size": 14}),
         ("Full ranked list of picks", {"size": 12, "color": TEXT_MID}),
         ("$29 / month",               {"bold": True, "color": POS, "size": 14, "align": PP_ALIGN.RIGHT})],
    ]
    add_table(s, lx + CARD_PAD, top_y + HDR_H + 0.1,
              col_w - 2 * CARD_PAD, top_h - HDR_H - 0.3,
              pricing_rows,
              col_widths=[1.8, col_w - 2 * CARD_PAD - 1.8 - 1.7, 1.7],
              row_heights=[0.35, 0.42, 0.42, 0.42])

    # Unit economics — native table
    add_card(s, rx, top_y, col_w, top_h,
             fill=BG_SURFACE, border=ACCENT,
             header_label="UNIT ECONOMICS", header_color=ACCENT)
    ue_rows = [
        [("Infrastructure cost",          {"size": 13, "color": TEXT_MID, "bold": False}),
         ("Under $50 / month",            {"size": 14, "color": POS, "bold": True, "align": PP_ALIGN.RIGHT})],
        [("Cost per additional user",     {"size": 13, "color": TEXT_MID, "bold": False}),
         ("Near zero",                    {"size": 14, "color": POS, "bold": True, "align": PP_ALIGN.RIGHT})],
        [("Gross margin at scale",        {"size": 13, "color": TEXT_MID, "bold": False}),
         ("95% or higher",                {"size": 14, "color": POS, "bold": True, "align": PP_ALIGN.RIGHT})],
        [("Capacity without rebuilding",  {"size": 13, "color": TEXT_MID, "bold": False}),
         ("10,000+ users",                {"size": 14, "color": POS, "bold": True, "align": PP_ALIGN.RIGHT})],
    ]
    add_table(s, rx + CARD_PAD, top_y + HDR_H + 0.15,
              col_w - 2 * CARD_PAD, top_h - HDR_H - 0.3,
              ue_rows, header_row=False,
              col_widths=[col_w - 2 * CARD_PAD - 2.0, 2.0],
              row_heights=[0.37, 0.37, 0.37, 0.37])

    # Bottom row — roadmap + risks
    bot_y = 5.0
    bot_h = 1.95
    bullets_font_size = 12
    def bullet_card(x, title, title_color, items):
        add_card(s, x, bot_y, col_w, bot_h, fill=BG_CARD, border=TEXT_DIM,
                 header_label=title, header_color=title_color)
        # Each bullet on its own line, vertically-distributed
        inner_h = bot_h - HDR_H - 0.2
        line_h = inner_h / len(items)
        for i, it in enumerate(items):
            add_text(s, x + CARD_PAD,
                     bot_y + HDR_H + 0.1 + line_h * i,
                     col_w - 2 * CARD_PAD, line_h,
                     "•  " + it,
                     size=bullets_font_size, color=TEXT_HI,
                     anchor=MSO_ANCHOR.MIDDLE)

    bullet_card(lx, "NEXT 90 DAYS", WARN, [
        "Public live scoreboard — weekly realized returns",
        "Options-market signals and earnings calendar",
        "Short-interest history as a capital-commitment signal",
        "Twelve-month portfolio simulator for prospects",
    ])
    bullet_card(rx, "WHAT WE WATCH", WARN, [
        "Past performance never guarantees the future",
        "Regime shifts kill signals — we monitor and retire them",
        "Fully systematic, by design — no human overlay",
        "Dependent on public US market data",
    ])

    add_page_number(s, 3)


# ── Build + thumbnail ─────────────────────────────────────────────
def build_full():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide_one(prs)
    slide_two(prs)
    slide_three(prs)
    prs.save(str(OUT))
    return OUT


def build_single(fn, suffix):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    fn(prs)
    path = Path(__file__).parent / f"tmp_{suffix}.pptx"
    prs.save(str(path))
    return path


def thumbnail(pptx_path: Path) -> Path:
    subprocess.run(["qlmanage", "-t", "-s", "1800", "-o", "/tmp", str(pptx_path)],
                    capture_output=True, timeout=60)
    return Path("/tmp") / (pptx_path.name + ".png")


def main() -> None:
    full = build_full()
    print(f"wrote {full} ({full.stat().st_size} bytes)")
    for i, fn in enumerate([slide_one, slide_two, slide_three], 1):
        p = build_single(fn, f"v4slide{i}")
        png = thumbnail(p)
        dest = Path("/tmp") / f"s-tool_deck_v4_slide{i}.png"
        if png.exists():
            shutil.copy(png, dest)
            p.unlink(missing_ok=True)
            print(f"slide {i}: {dest}")
        else:
            print(f"slide {i}: MISSING (pptx saved; qlmanage failed)")


if __name__ == "__main__":
    main()
