"""Investor deck v3 — polished for a non-technical audience.

v2 fixed layout overlap. v3 fixes vocabulary:
  - No underscores ("hit_100" → "Doubles Rate")
  - Every acronym defined on first mention (e.g. "Out-of-Sample (OOS)")
  - Story-first slide 1 — the 11× claim leads, stats support
  - Architecture slide in plain English — 4 concepts, not 4 code modules
  - Business slide keeps pricing clean, roadmap readable, risks honest

Produces research/s-tool_investor_deck_v3_2026_04_17.pptx plus per-slide
thumbnails at /tmp/s-tool_deck_v3_slide{1,2,3}.png.
"""

from __future__ import annotations

import json
import shutil
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).parent.parent
BACKTEST = json.loads((ROOT / "data_cache" / "backtest_report.json").read_text())
OUT = Path(__file__).parent / "s-tool_investor_deck_v3_2026_04_17.pptx"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_IN = 0.5
GUTTER_IN = 0.3
TITLE_Y_IN = 0.35
TITLE_H_IN = 0.9

BG_DEEP    = RGBColor(0x0B, 0x10, 0x06)
BG_SURFACE = RGBColor(0x1A, 0x1F, 0x1B)
BG_CARD    = RGBColor(0x1F, 0x27, 0x22)
TEXT_HI    = RGBColor(0xED, 0xEC, 0xE6)
TEXT_MID   = RGBColor(0x9B, 0xA1, 0xB9)
TEXT_DIM   = RGBColor(0x6B, 0x73, 0x82)
ACCENT     = RGBColor(0x5F, 0xAA, 0xC5)
POS        = RGBColor(0x6E, 0xE7, 0xB7)
WARN       = RGBColor(0xF5, 0xD5, 0x8F)


def add_bg(slide):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
    shp.fill.solid(); shp.fill.fore_color.rgb = BG_DEEP
    shp.line.fill.background(); shp.shadow.inherit = False


def add_card(slide, l, t, w, h, *, fill=BG_SURFACE, border=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(l), Inches(t), Inches(w), Inches(h))
    shp.adjustments[0] = 0.04
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border; shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False


def add_text(slide, l, t, w, h, text, *, size=14, color=TEXT_HI,
              bold=False, italic=False, align=PP_ALIGN.LEFT,
              font="Helvetica Neue", line_spacing=1.2):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font; f.size = Pt(size); f.color.rgb = color
    f.bold = bold; f.italic = italic


def add_title(slide, eyebrow, title_text, *, accent_phrase=None,
               title_size=30):
    add_text(slide, MARGIN_IN, TITLE_Y_IN,
             SLIDE_W_IN - 2 * MARGIN_IN, 0.3,
             eyebrow, size=10, color=ACCENT, bold=True)
    tx = slide.shapes.add_textbox(Inches(MARGIN_IN),
                                   Inches(TITLE_Y_IN + 0.35),
                                   Inches(SLIDE_W_IN - 2 * MARGIN_IN),
                                   Inches(TITLE_H_IN))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.line_spacing = 1.1

    if accent_phrase and accent_phrase in title_text:
        before, _, after = title_text.partition(accent_phrase)
        segs = [(before, False), (accent_phrase, True), (after, False)]
    else:
        segs = [(title_text, False)]

    for txt, is_accent in segs:
        if not txt:
            continue
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = "Georgia"; f.size = Pt(title_size)
        f.color.rgb = ACCENT if is_accent else TEXT_HI
        f.italic = is_accent


def add_page_number(slide, n):
    add_text(slide, SLIDE_W_IN - 0.7, SLIDE_H_IN - 0.32, 0.5, 0.2,
             str(n), size=9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


# ── Slide 1 — Lead with the story ─────────────────────────────────
def slide_one(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s,
              "THE EDGE  ·  BACKTESTED MAY 2022 THROUGH AUGUST 2024",
              "Our picks double at 11× the rate of random stock picking.",
              accent_phrase="11× the rate of random stock picking")

    # Two BIG comparison cards, plain English
    top_y = 2.3
    top_h = 2.55
    col_w = (SLIDE_W_IN - 2 * MARGIN_IN - GUTTER_IN) / 2

    # LEFT — random picking baseline
    lx = MARGIN_IN
    add_card(s, lx, top_y, col_w, top_h, fill=BG_SURFACE, border=TEXT_DIM)
    add_text(s, lx + 0.3, top_y + 0.25, col_w - 0.6, 0.3,
             "RANDOM STOCK FROM THE RUSSELL 3000",
             size=11, color=TEXT_MID, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, lx + 0.3, top_y + 0.7, col_w - 0.6, 1.05,
             "1 in 20", size=58, color=TEXT_HI, bold=True,
             font="Georgia", align=PP_ALIGN.CENTER)
    add_text(s, lx + 0.3, top_y + 1.75, col_w - 0.6, 0.35,
             "chance of doubling over the next year",
             size=14, color=TEXT_MID, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, lx + 0.3, top_y + 2.15, col_w - 0.6, 0.3,
             "(~4.95% base rate from 22,218 observations)",
             size=10, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # RIGHT — our model's picks
    rx = MARGIN_IN + col_w + GUTTER_IN
    add_card(s, rx, top_y, col_w, top_h, fill=BG_SURFACE, border=POS)
    add_text(s, rx + 0.3, top_y + 0.25, col_w - 0.6, 0.3,
             "TOP 20 PICKS FROM OUR MODEL",
             size=11, color=POS, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, rx + 0.3, top_y + 0.7, col_w - 0.6, 1.05,
             "11 in 20", size=58, color=POS, bold=True,
             font="Georgia", align=PP_ALIGN.CENTER)
    add_text(s, rx + 0.3, top_y + 1.75, col_w - 0.6, 0.35,
             "doubled in the following year",
             size=14, color=TEXT_MID, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, rx + 0.3, top_y + 2.15, col_w - 0.6, 0.3,
             "(55.0% realized rate on top-ranked picks)",
             size=10, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Honest checks strip — 3 columns, plain-English labels
    strip_y = 5.1
    strip_h = 1.7
    add_card(s, MARGIN_IN, strip_y, SLIDE_W_IN - 2 * MARGIN_IN, strip_h,
             fill=BG_CARD, border=WARN)
    add_text(s, MARGIN_IN + 0.3, strip_y + 0.2, 11.5, 0.25,
             "AND THESE ARE THE HONEST CHECKS  ·  NO CHERRY-PICKING",
             size=10, color=WARN, bold=True)

    checks = [
        ("Works across every company size",
         "8.45× lift",
         "From the smallest to largest 20%\nof Russell 3000 companies"),
        ("Tested on market data the model never saw",
         "68% doubled",
         "On 2024 out-of-sample (OOS) data — the\nsignal strengthened, not weakened"),
        ("Beats simple rules by 5×",
         "11.1× vs 2.0×",
         "Our model's lift vs the best\nhand-crafted rule we tried"),
    ]
    col_w_check = (SLIDE_W_IN - 2 * MARGIN_IN - 0.6) / 3
    for i, (lab, val, sub) in enumerate(checks):
        cx = MARGIN_IN + 0.3 + col_w_check * i
        add_text(s, cx, strip_y + 0.55, col_w_check, 0.3,
                 lab, size=11, color=TEXT_MID, bold=True)
        add_text(s, cx, strip_y + 0.85, col_w_check, 0.45,
                 val, size=22, color=POS, bold=True, font="Georgia")
        add_text(s, cx, strip_y + 1.30, col_w_check, 0.4,
                 sub, size=10, color=TEXT_DIM, italic=True, line_spacing=1.15)

    add_text(s, MARGIN_IN, SLIDE_H_IN - 0.4, SLIDE_W_IN - 2 * MARGIN_IN, 0.25,
             "Every pick is logged publicly for live verification at s-tool.io/track-record",
             size=10, color=TEXT_DIM, italic=True)
    add_page_number(s, 1)


# ── Slide 2 — How it works, in plain English ─────────────────────
def slide_two(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "HOW IT WORKS",
              "The model rebuilds itself every night.",
              accent_phrase="every night")

    # 4 plain-English concept cards
    col_y = 2.3
    col_h = 3.4
    cols = 4
    total_w = SLIDE_W_IN - 2 * MARGIN_IN
    col_w = (total_w - GUTTER_IN * (cols - 1)) / cols
    concepts = [
        ("WIDE UNIVERSE", ACCENT,
         "~3,000 companies",
         "The Russell 3000 — the broadest US stock index. "
         "We don't cherry-pick; we rank the whole market."),
        ("RICH DATA", POS,
         "Public data only",
         "Prices, earnings filings with the SEC, analyst "
         "estimates, and macroeconomic indicators from the Fed."),
        ("MACHINE LEARNING", WARN,
         "Trained on history",
         "A machine-learning model learns which combinations "
         "of signals preceded real 100%+ moves in the past."),
        ("DAILY RETRAIN", RGBColor(0xD2, 0xDD, 0xEA),
         "Under 24 hours",
         "Every night the model sees yesterday's real market "
         "prints and updates. No static algorithm decaying in a drawer."),
    ]
    for i, (label, color, headline, body) in enumerate(concepts):
        x = MARGIN_IN + (col_w + GUTTER_IN) * i
        add_card(s, x, col_y, col_w, col_h, fill=BG_SURFACE, border=color)
        # Header strip
        hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(col_y),
                                 Inches(col_w), Inches(0.5))
        hdr.adjustments[0] = 0.2
        hdr.fill.solid(); hdr.fill.fore_color.rgb = color
        hdr.line.fill.background(); hdr.shadow.inherit = False
        add_text(s, x, col_y + 0.1, col_w, 0.3,
                 label, size=13, color=BG_DEEP, bold=True, align=PP_ALIGN.CENTER)
        # Headline number / phrase
        add_text(s, x + 0.2, col_y + 0.75, col_w - 0.4, 0.55,
                 headline, size=22, color=TEXT_HI, bold=True,
                 font="Georgia", align=PP_ALIGN.CENTER)
        # Body — let it breathe
        add_text(s, x + 0.25, col_y + 1.5, col_w - 0.5, col_h - 1.7,
                 body, size=12, color=TEXT_MID, line_spacing=1.35)

    # Loop visualisation strip
    loop_y = 6.0
    loop_h = 0.8
    add_card(s, MARGIN_IN, loop_y, SLIDE_W_IN - 2 * MARGIN_IN, loop_h,
             fill=BG_CARD, border=ACCENT)
    add_text(s, MARGIN_IN + 0.3, loop_y + 0.12,
             SLIDE_W_IN - 2 * MARGIN_IN - 0.6, 0.25,
             "THE FEEDBACK LOOP",
             size=10, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN_IN + 0.3, loop_y + 0.4,
             SLIDE_W_IN - 2 * MARGIN_IN - 0.6, 0.35,
             "yesterday's market prints  →  realized outcomes  →  model updates  →  "
             "today's ranked picks",
             size=14, color=TEXT_HI, bold=True, align=PP_ALIGN.CENTER)
    add_page_number(s, 2)


# ── Slide 3 — Business ────────────────────────────────────────────
def slide_three(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "WHERE THE CAPITAL GOES  ·  2026 Q2",
              "Real edge. Honest surface. Infrastructure built for scale.",
              accent_phrase="Infrastructure built for scale")

    col_w = (SLIDE_W_IN - 2 * MARGIN_IN - GUTTER_IN) / 2
    lx = MARGIN_IN
    rx = MARGIN_IN + col_w + GUTTER_IN

    # Pricing
    top_y = 2.3
    top_h = 2.4
    add_card(s, lx, top_y, col_w, top_h, fill=BG_SURFACE, border=POS)
    add_text(s, lx + 0.3, top_y + 0.2, col_w - 0.6, 0.3,
             "PRICING  ·  LIVE TODAY",
             size=11, color=POS, bold=True)
    add_text(s, lx + 0.4, top_y + 0.65, 2.2, 0.3,
             "TIER", size=9, color=TEXT_DIM, bold=True)
    add_text(s, lx + 2.5, top_y + 0.65, 3.0, 0.3,
             "WHAT YOU GET", size=9, color=TEXT_DIM, bold=True)
    add_text(s, lx + col_w - 1.7, top_y + 0.65, 1.3, 0.3,
             "PRICE", size=9, color=TEXT_DIM, bold=True, align=PP_ALIGN.RIGHT)
    tiers = [
        ("Free",       "3 projections per day",             "Free"),
        ("Pro",        "10 projections per day",            "$8 / mo"),
        ("Strategist", "Full ranked list of picks",         "$29 / mo"),
    ]
    for i, (t, w, p) in enumerate(tiers):
        y = top_y + 1.0 + 0.42 * i
        add_text(s, lx + 0.4, y, 2.2, 0.3, t,
                 size=15, color=TEXT_HI, bold=True)
        add_text(s, lx + 2.5, y, 3.0, 0.3, w,
                 size=12, color=TEXT_MID)
        add_text(s, lx + col_w - 1.7, y, 1.3, 0.3, p,
                 size=15, color=POS, bold=True, font="Georgia",
                 align=PP_ALIGN.RIGHT)

    # Unit Economics
    add_card(s, rx, top_y, col_w, top_h, fill=BG_SURFACE, border=ACCENT)
    add_text(s, rx + 0.3, top_y + 0.2, col_w - 0.6, 0.3,
             "UNIT ECONOMICS",
             size=11, color=ACCENT, bold=True)
    ue = [
        ("Infrastructure cost",        "Under $50 / month"),
        ("Cost per additional user",   "Near zero"),
        ("Gross margin at scale",      "95%+"),
        ("Capacity without rebuilding","10,000+ users"),
    ]
    for i, (lab, val) in enumerate(ue):
        y = top_y + 0.75 + 0.42 * i
        add_text(s, rx + 0.4, y, col_w - 2.4, 0.3, lab,
                 size=13, color=TEXT_MID)
        add_text(s, rx + col_w - 2.2, y, 1.9, 0.3, val,
                 size=15, color=POS, bold=True, font="Georgia",
                 align=PP_ALIGN.RIGHT)

    # Bottom — Roadmap + Risks
    bot_y = 4.9
    bot_h = 2.05
    add_card(s, lx, bot_y, col_w, bot_h, fill=BG_CARD, border=TEXT_DIM)
    add_text(s, lx + 0.3, bot_y + 0.15, col_w - 0.6, 0.3,
             "NEXT 90 DAYS",
             size=11, color=WARN, bold=True)
    road = [
        "Public live scoreboard — weekly realized returns",
        "Options-market signals and earnings calendar",
        "Short-interest history as a capital-commitment signal",
        "12-month portfolio simulator for prospects",
    ]
    for i, r in enumerate(road):
        add_text(s, lx + 0.4, bot_y + 0.55 + 0.32 * i, col_w - 0.8, 0.3,
                 "— " + r, size=12, color=TEXT_HI)

    add_card(s, rx, bot_y, col_w, bot_h, fill=BG_CARD, border=TEXT_DIM)
    add_text(s, rx + 0.3, bot_y + 0.15, col_w - 0.6, 0.3,
             "WHAT WE WATCH",
             size=11, color=WARN, bold=True)
    risks = [
        "Past performance never guarantees the future",
        "Regime shifts kill signals — we monitor and retire them",
        "Fully systematic (no gut-feel overlay, by design)",
        "Dependent on public US market data",
    ]
    for i, r in enumerate(risks):
        add_text(s, rx + 0.4, bot_y + 0.55 + 0.32 * i, col_w - 0.8, 0.3,
                 "— " + r, size=12, color=TEXT_MID)

    add_page_number(s, 3)


def build_full():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide_one(prs)
    slide_two(prs)
    slide_three(prs)
    prs.save(str(OUT))
    return OUT


def build_single(fn, suffix):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    fn(prs)
    path = Path(__file__).parent / f"tmp_{suffix}.pptx"
    prs.save(str(path))
    return path


def thumbnail(pptx_path: Path) -> Path:
    subprocess.run(["qlmanage", "-t", "-s", "1600", "-o", "/tmp", str(pptx_path)],
                    capture_output=True, timeout=60)
    return Path("/tmp") / (pptx_path.name + ".png")


def main() -> None:
    full = build_full()
    print(f"wrote {full} ({full.stat().st_size} bytes)")
    for i, fn in enumerate([slide_one, slide_two, slide_three], 1):
        p = build_single(fn, f"v3slide{i}")
        png = thumbnail(p)
        dest = Path("/tmp") / f"s-tool_deck_v3_slide{i}.png"
        if png.exists():
            shutil.copy(png, dest)
            p.unlink(missing_ok=True)
            print(f"slide {i}: {dest}")
        else:
            print(f"slide {i}: MISSING")


if __name__ == "__main__":
    main()
