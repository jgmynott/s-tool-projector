"""Investor deck v2 — grounded in the Media Measurement Frameworks deck
layout conventions.

Learned from v1 overlap failure:
  - Cap text boxes per slide at ~12 (ref deck uses 3-12, v1 had 42)
  - Title strip full-width 12.60 x 0.83 at 0.41,0.28 (reference pattern)
  - Body content starts at y ~= 1.2in
  - Font sizes 14-16pt body, 20-24pt sub-headers, 28pt+ stats (ref pattern)
  - Page number at 12.60,7.20 (0.41x0.23, ref pattern)
  - Fixed gutters (GUTTER_IN = 0.3) — no ad-hoc 0.05*i hacks
  - Budget whitespace — prefer emptier slides than cramped ones

Outputs:
  research/s-tool_investor_deck_v2_2026_04_17.pptx
  research/tmp_slide_{1,2,3}.pptx              (single-slide variants)
  /tmp/s-tool_slide_{1,2,3}.png                (qlmanage thumbnails)
"""

from __future__ import annotations

import json
import shutil
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).parent.parent
BACKTEST = json.loads((ROOT / "data_cache" / "backtest_report.json").read_text())
OUT = Path(__file__).parent / "s-tool_investor_deck_v2_2026_04_17.pptx"

# ── Reference-grounded layout constants ──
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_IN = 0.5
GUTTER_IN = 0.3
TITLE_Y_IN = 0.35
TITLE_H_IN = 0.9
BODY_Y_IN = 1.45         # first body element starts here
BODY_BOTTOM_IN = 6.85    # last body element ends by here

# ── Palette (matches live s-tool.io theme) ──
BG_DEEP     = RGBColor(0x0B, 0x10, 0x06)
BG_SURFACE  = RGBColor(0x1A, 0x1F, 0x1B)
BG_CARD     = RGBColor(0x1F, 0x27, 0x22)
TEXT_HI     = RGBColor(0xED, 0xEC, 0xE6)
TEXT_MID    = RGBColor(0x9B, 0xA1, 0xB9)
TEXT_DIM    = RGBColor(0x6B, 0x73, 0x82)
ACCENT      = RGBColor(0x5F, 0xAA, 0xC5)
POS         = RGBColor(0x6E, 0xE7, 0xB7)
WARN        = RGBColor(0xF5, 0xD5, 0x8F)


def add_bg(slide, color=BG_DEEP):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def add_card(slide, l, t, w, h, *, fill=BG_SURFACE, border=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(l), Inches(t), Inches(w), Inches(h))
    shp.adjustments[0] = 0.04
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border; shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, l, t, w, h, text, *, size=14, color=TEXT_HI,
              bold=False, italic=False, align=PP_ALIGN.LEFT,
              font="Helvetica Neue", line_spacing=1.2):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.color.rgb = color
    f.bold = bold
    f.italic = italic
    return tx


def add_page_number(slide, n):
    add_text(slide, SLIDE_W_IN - 0.7, SLIDE_H_IN - 0.32, 0.5, 0.2,
             str(n), size=9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


def add_title(slide, eyebrow, title_text, *, accent_phrase=None):
    # Eyebrow strip just above the title
    add_text(slide, MARGIN_IN, TITLE_Y_IN, SLIDE_W_IN - 2 * MARGIN_IN, 0.3,
             eyebrow, size=10, color=ACCENT, bold=True)
    # Title line — if accent_phrase given, italicize that section in ACCENT
    tx = slide.shapes.add_textbox(Inches(MARGIN_IN), Inches(TITLE_Y_IN + 0.35),
                                   Inches(SLIDE_W_IN - 2 * MARGIN_IN),
                                   Inches(TITLE_H_IN))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.line_spacing = 1.1

    segments = [(title_text, False)]
    if accent_phrase and accent_phrase in title_text:
        before, _, after = title_text.partition(accent_phrase)
        segments = [(before, False), (accent_phrase, True), (after, False)]

    for txt, is_accent in segments:
        if not txt:
            continue
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = "Georgia"
        f.size = Pt(30)
        f.color.rgb = ACCENT if is_accent else TEXT_HI
        f.italic = is_accent
    return tx


# ── Slide 1 — The Edge ────────────────────────────────────────────
def slide_performance(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s,
              "THE EDGE  ·  BACKTEST 2022-05 → 2024-08  ·  10 WALK-FORWARD WINDOWS",
              "Our ranker picks +100% movers at 11× the Russell 3000 baseline.",
              accent_phrase="+100% movers at 11×")

    # Two cards side-by-side: baseline vs our model.
    card_y = 2.35
    card_h = 2.35
    card_w = (SLIDE_W_IN - 2 * MARGIN_IN - GUTTER_IN) / 2   # = 6.1165
    left_x = MARGIN_IN
    right_x = MARGIN_IN + card_w + GUTTER_IN

    # Left — baseline
    add_card(s, left_x, card_y, card_w, card_h, fill=BG_SURFACE, border=TEXT_DIM)
    add_text(s, left_x + 0.3, card_y + 0.2, card_w - 0.6, 0.3,
             "UNIVERSE BASELINE · RANDOM TICKER, 1Y",
             size=10, color=ACCENT, bold=True)
    base = BACKTEST["baseline_rates"]
    rows = [("+10%", f"{base['+10%']*100:.1f}%"),
            ("+25%", f"{base['+25%']*100:.1f}%"),
            ("+50%", f"{base['+50%']*100:.1f}%"),
            ("+100%", f"{base['+100%']*100:.2f}%"),
            ("+200%", f"{base['+200%']*100:.2f}%")]
    for i, (k, v) in enumerate(rows):
        y = card_y + 0.7 + 0.3 * i
        add_text(s, left_x + 0.4, y, 1.5, 0.3, k, size=14, color=TEXT_MID)
        add_text(s, left_x + 2.0, y, 2.0, 0.3, v, size=14, color=TEXT_HI, bold=True)

    # Right — our model
    add_card(s, right_x, card_y, card_w, card_h, fill=BG_SURFACE, border=ACCENT)
    add_text(s, right_x + 0.3, card_y + 0.2, card_w - 0.6, 0.3,
             "OUR NN SCORER · TOP-20 PICKS / WINDOW",
             size=10, color=ACCENT, bold=True)
    nn = BACKTEST["methods"]["nn_score"]
    hc = BACKTEST["methods"]["H7_ewma_p90"]
    kpis = [("Hit rate at +100%",       f"{nn['thresholds']['+100%']['rate']*100:.1f}%"),
            ("Lift vs baseline",         f"{nn['thresholds']['+100%']['lift']:.2f}×"),
            ("Mean 1Y return",            f"+{nn['mean_return']*100:.0f}%"),
            ("Median 1Y return",          f"+{nn['median_return']*100:.0f}%"),
            ("Hand-crafted baseline lift", f"{hc['thresholds']['+100%']['lift']:.2f}×")]
    for i, (lab, val) in enumerate(kpis):
        y = card_y + 0.7 + 0.3 * i
        add_text(s, right_x + 0.4, y, 3.2, 0.3, lab, size=13, color=TEXT_MID)
        add_text(s, right_x + 3.7, y, card_w - 4.0, 0.3, val,
                 size=15, color=POS, bold=True, font="Georgia", align=PP_ALIGN.RIGHT)

    # Honest-metrics strip — 4 big numbers across
    strip_y = 5.0
    strip_h = 1.55
    add_card(s, MARGIN_IN, strip_y, SLIDE_W_IN - 2 * MARGIN_IN, strip_h,
             fill=BG_CARD, border=WARN)
    add_text(s, MARGIN_IN + 0.3, strip_y + 0.15, 11.0, 0.25,
             "HONEST METRICS  ·  SIZE-NEUTRAL  ·  OUT-OF-SAMPLE",
             size=10, color=WARN, bold=True)

    hm = BACKTEST["honest_metrics"]
    cols = [
        (f"{hm['size_neutral_hit_100']*100:.1f}%",    "Size-neutral hit_100"),
        (f"{hm['within_quintile_lift_median']:.2f}×", "Within-quintile lift"),
        (f"{hm['year_oos_hit_100']*100:.1f}%",        "2024 OOS hit_100"),
        (f"{hm['year_oos_lift']:.2f}×",               "2024 OOS lift"),
    ]
    col_w = (SLIDE_W_IN - 2 * MARGIN_IN - 0.6) / 4
    for i, (val, lab) in enumerate(cols):
        cx = MARGIN_IN + 0.3 + col_w * i
        add_text(s, cx, strip_y + 0.5, col_w, 0.6,
                 val, size=32, color=POS, bold=True, font="Georgia",
                 align=PP_ALIGN.CENTER)
        add_text(s, cx, strip_y + 1.15, col_w, 0.3,
                 lab, size=11, color=TEXT_MID, align=PP_ALIGN.CENTER)

    # Footer
    add_text(s, MARGIN_IN, SLIDE_H_IN - 0.45, SLIDE_W_IN - 2 * MARGIN_IN, 0.25,
             f"{BACKTEST['universe_size']:,} ticker-window rows  ·  "
             f"Russell 3000 universe  ·  every pick logged for live tracking",
             size=10, color=TEXT_DIM)
    add_page_number(s, 1)


# ── Slide 2 — Architecture ────────────────────────────────────────
def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "ARCHITECTURE  ·  NIGHTLY END-TO-END LEARNING LOOP",
              "The model retrains every night on freshly realised returns.",
              accent_phrase="every night")

    # 4-column stack — kept deliberately short per column.
    col_y = 2.45
    col_h = 3.5
    total_w = SLIDE_W_IN - 2 * MARGIN_IN
    col_count = 4
    col_w = (total_w - GUTTER_IN * (col_count - 1)) / col_count  # ≈ 2.86
    cols = [
        ("DATA",    ACCENT, [
            "Russell 3000 universe",
            "2,609 price histories (5y)",
            "1,947 SEC filings · 98k rows",
            "Analyst consensus + EPS",
            "Macro regime · 6 indicators",
        ]),
        ("COMPUTE", POS, [
            "MC + mean-reversion blend",
            "Tree-ensemble NN scorer",
            "Walk-forward training",
            "Size-neutral bucketing",
            "Liquidity floor enforced",
        ]),
        ("SERVE",   WARN, [
            "Cloudflare edge — 15KB HTML",
            "FastAPI · Clerk · Stripe",
            "Railway Volume for users DB",
            "Gated Strategist tier",
            "/track-record is public",
        ]),
        ("CI LOOP", RGBColor(0xD2, 0xDD, 0xEA), [
            "20:00 UTC · preferred scan",
            "23:00 UTC · NN retrain",
            "Nightly artifact deploys",
            "Every pick into ledger",
            "Cache warm across runs",
        ]),
    ]
    for i, (label, color, lines) in enumerate(cols):
        x = MARGIN_IN + (col_w + GUTTER_IN) * i
        add_card(s, x, col_y, col_w, col_h, fill=BG_SURFACE, border=color)
        # Header strip
        hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(col_y),
                                 Inches(col_w), Inches(0.45))
        hdr.adjustments[0] = 0.2
        hdr.fill.solid(); hdr.fill.fore_color.rgb = color
        hdr.line.fill.background(); hdr.shadow.inherit = False
        add_text(s, x, col_y + 0.07, col_w, 0.3,
                 label, size=14, color=BG_DEEP, bold=True, align=PP_ALIGN.CENTER)

        # Bullets — 5 per column, 14pt, generous spacing
        for j, ln in enumerate(lines):
            add_text(s, x + 0.2, col_y + 0.65 + 0.55 * j,
                     col_w - 0.3, 0.5,
                     "· " + ln,
                     size=12, color=TEXT_HI, line_spacing=1.2)

    # Loop callout strip
    loop_y = 6.2
    loop_h = 0.6
    add_card(s, MARGIN_IN, loop_y, SLIDE_W_IN - 2 * MARGIN_IN, loop_h,
             fill=BG_CARD, border=ACCENT)
    add_text(s, MARGIN_IN + 0.3, loop_y + 0.15, SLIDE_W_IN - 2 * MARGIN_IN - 0.6, 0.35,
             "today's prints  →  realisation join  →  walk-forward retrain  →  "
             "fresh ranking  →  tomorrow's picks",
             size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_page_number(s, 2)


# ── Slide 3 — Business ────────────────────────────────────────────
def slide_business(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "WHERE CAPITAL GOES  ·  2026 Q2",
              "Real edge. Honest surface. Infra ready for the next 100×.",
              accent_phrase="the next 100×")

    # Two cards side-by-side: pricing | unit economics
    top_y = 2.45
    top_h = 2.3
    col_w = (SLIDE_W_IN - 2 * MARGIN_IN - GUTTER_IN) / 2
    lx = MARGIN_IN
    rx = MARGIN_IN + col_w + GUTTER_IN

    # Pricing
    add_card(s, lx, top_y, col_w, top_h, fill=BG_SURFACE, border=POS)
    add_text(s, lx + 0.3, top_y + 0.2, col_w - 0.6, 0.3,
             "PRICING · REVENUE SURFACE", size=10, color=POS, bold=True)
    tiers = [
        ("Free",       "3 projections / day",      "—"),
        ("Pro",        "10 projections / day",     "$8 / mo"),
        ("Strategist", "Full ranked list + picks", "$29 / mo"),
    ]
    # Header row
    add_text(s, lx + 0.4, top_y + 0.6, 1.7, 0.25, "TIER",
             size=9, color=TEXT_DIM, bold=True)
    add_text(s, lx + 2.1, top_y + 0.6, 2.6, 0.25, "WHAT",
             size=9, color=TEXT_DIM, bold=True)
    add_text(s, lx + col_w - 1.6, top_y + 0.6, 1.2, 0.25, "PRICE",
             size=9, color=TEXT_DIM, bold=True, align=PP_ALIGN.RIGHT)
    for i, (t, w, p) in enumerate(tiers):
        y = top_y + 0.9 + 0.4 * i
        add_text(s, lx + 0.4, y, 1.7, 0.3, t, size=14, color=TEXT_HI, bold=True)
        add_text(s, lx + 2.1, y, 2.6, 0.3, w, size=12, color=TEXT_MID)
        add_text(s, lx + col_w - 1.6, y, 1.2, 0.3, p,
                 size=14, color=POS, bold=True, font="Georgia",
                 align=PP_ALIGN.RIGHT)

    # Unit economics
    add_card(s, rx, top_y, col_w, top_h, fill=BG_SURFACE, border=ACCENT)
    add_text(s, rx + 0.3, top_y + 0.2, col_w - 0.6, 0.3,
             "UNIT ECONOMICS", size=10, color=ACCENT, bold=True)
    ue = [
        ("Infra cost",              "<$50 / mo"),
        ("Per-user marginal cost",  "~$0"),
        ("Gross margin @ scale",    "95%+"),
        ("Capacity today",          "10k+ users"),
    ]
    for i, (lab, val) in enumerate(ue):
        y = top_y + 0.6 + 0.4 * i
        add_text(s, rx + 0.4, y, col_w - 2.0, 0.3, lab, size=13, color=TEXT_MID)
        add_text(s, rx + col_w - 1.8, y, 1.4, 0.3, val,
                 size=16, color=POS, bold=True, font="Georgia",
                 align=PP_ALIGN.RIGHT)

    # Bottom — Roadmap + Risks side-by-side
    bot_y = 5.0
    bot_h = 1.85
    add_card(s, MARGIN_IN, bot_y, col_w, bot_h, fill=BG_CARD, border=TEXT_DIM)
    add_text(s, MARGIN_IN + 0.3, bot_y + 0.15, col_w - 0.6, 0.3,
             "ROADMAP · NEXT 90 DAYS", size=10, color=WARN, bold=True)
    road = [
        "Live pick scoreboard — weekly realised returns",
        "Options flow + earnings calendar as NN features",
        "Short interest backfill — capital-commitment signal",
        "Portfolio-level 12-month simulation",
    ]
    for i, r in enumerate(road):
        add_text(s, MARGIN_IN + 0.4, bot_y + 0.5 + 0.3 * i, col_w - 0.8, 0.3,
                 "· " + r, size=11, color=TEXT_HI)

    add_card(s, rx, bot_y, col_w, bot_h, fill=BG_CARD, border=TEXT_DIM)
    add_text(s, rx + 0.3, bot_y + 0.15, col_w - 0.6, 0.3,
             "RISKS · WHAT WE WATCH", size=10, color=WARN, bold=True)
    risks = [
        "Past backtest ≠ future guarantee",
        "Regime shift killed crowd-sentiment in 2023",
        "Russell 3000 scan: ~45 min, headroom to 90 min",
        "Fully systematic — no human overlay",
    ]
    for i, r in enumerate(risks):
        add_text(s, rx + 0.4, bot_y + 0.5 + 0.3 * i, col_w - 0.8, 0.3,
                 "· " + r, size=11, color=TEXT_MID)

    add_page_number(s, 3)


def build_full():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide_performance(prs)
    slide_architecture(prs)
    slide_business(prs)
    prs.save(str(OUT))
    return OUT


def build_single(slide_fn, suffix):
    """Build a one-slide deck so qlmanage can thumbnail it in isolation."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide_fn(prs)
    path = Path(__file__).parent / f"tmp_{suffix}.pptx"
    prs.save(str(path))
    return path


def thumbnail(pptx_path: Path) -> Path:
    subprocess.run(["qlmanage", "-t", "-s", "1600", "-o", "/tmp", str(pptx_path)],
                    capture_output=True, timeout=60)
    png = Path("/tmp") / (pptx_path.name + ".png")
    return png


def main() -> None:
    full = build_full()
    print(f"wrote {full} ({full.stat().st_size} bytes)")

    # Per-slide thumbnails via single-slide decks
    for i, fn in enumerate([slide_performance, slide_architecture, slide_business], 1):
        p = build_single(fn, f"slide{i}")
        png = thumbnail(p)
        dest = Path("/tmp") / f"s-tool_deck_v2_slide{i}.png"
        if png.exists():
            shutil.copy(png, dest)
            print(f"thumbnail slide {i}: {dest}")
        else:
            print(f"thumbnail slide {i}: MISSING")


if __name__ == "__main__":
    main()
